import os
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv

class SlideGenerator:
    def __init__(self, output_dir="output/"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def create_slide_deck(self, title, bullet_points, images=None):
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        bullet_slide_layout = prs.slide_layouts[1]

        # Title Slide
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "AI Generated Presentation"
        slide.placeholders[1].text = f"Topic: {title}"

        for i, point in enumerate(bullet_points):
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = f"Slide {i + 1}"
            content = slide.placeholders[1]
            content.text = point

            # ✅ Add image if available and path exists
            if images and i < len(images):
                image_path = images[i]
                if os.path.exists(image_path):
                    try:
                        # Adjust position and size as needed
                        slide.shapes.add_picture(image_path, Inches(0.5), Inches(3.5), width=Inches(8))
                    except Exception as e:
                        print(f"Failed to add image '{image_path}' to slide {i + 1}: {e}")
                else:
                    print(f"Image not found: {image_path}")

        pptx_path = os.path.join(self.output_dir, "generated_presentation.pptx")
        prs.save(pptx_path)
        print(f"✅ Presentation saved to: {pptx_path}")

if __name__ == "__main__":
    # Example test
    sample_title = "Sample Presentation"
    sample_points = [
        "AI automates repetitive tasks.",
        "Data-driven decisions are enhanced.",
        "Documents converted to slides.",
        "Using GPT for summarization.",
        "LlamaIndex used for data handling."
    ]
    sample_images = [
        "output/images/sample_image_0.png",
        "output/images/sample_image_1.png",
        "output/images/sample_image_2.png",
        "output/images/sample_image_3.png",
        "output/images/sample_image_4.png"
    ]

    generator = SlideGenerator()
    generator.create_slide_deck(sample_title, sample_points, sample_images)
